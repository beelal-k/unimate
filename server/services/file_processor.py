# services/file_processor.py
# Text extraction from PDF, PPTX, DOCX, and plain text files

import io
import logging

logger = logging.getLogger(__name__)


def extract(file_bytes: bytes, mimetype: str) -> str:
    """Route to the correct extractor based on mimetype."""
    extractors = {
        "application/pdf": extract_pdf,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
        "text/plain": extract_txt,
        "text/html": extract_html,
    }

    extractor = extractors.get(mimetype)
    if not extractor:
        logger.warning("Unsupported mimetype for extraction: %s", mimetype)
        return ""

    try:
        return extractor(file_bytes)
    except Exception as e:
        logger.error("Failed to extract text from %s: %s", mimetype, e)
        return ""


def extract_all(attachments: list[dict]) -> str:
    """Extract text from all attachments and concatenate."""
    parts = []
    for attachment in attachments:
        filename = attachment.get("filename", "unknown")
        mimetype = attachment.get("mimetype", "")
        content = attachment.get("content", b"")

        text = extract(content, mimetype)
        if text and text.strip():
            parts.append(f"--- {filename} ---\n{text}")

    return "\n\n".join(parts)


def extract_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i}]\n{text.strip()}")

    return "\n\n".join(pages)


def extract_pptx(file_bytes: bytes) -> str:
    """
    Extract text from PPTX using python-pptx.
    Includes text frames, speaker notes, and table cells.
    Preserves slide structure with markers.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = [f"[Slide {i}]"]

        # Extract text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)

            # Extract table content
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_text.append(" | ".join(cells))

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"[Speaker Notes] {notes}")

        if len(slide_text) > 1:  # More than just the slide marker
            slides.append("\n".join(slide_text))

    return "\n\n".join(slides)


def extract_docx(file_bytes: bytes) -> str:
    """
    Extract text from DOCX using python-docx.
    Preserves heading structure and table content.
    """
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        # Preserve heading structure
        if paragraph.style and paragraph.style.name.startswith("Heading"):
            level = paragraph.style.name.replace("Heading ", "").strip()
            try:
                level_num = int(level)
                prefix = "#" * level_num
                parts.append(f"{prefix} {text}")
            except ValueError:
                parts.append(text)
        else:
            parts.append(text)

    # Extract table content
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
        if table_rows:
            parts.append("\n".join(table_rows))

    return "\n\n".join(parts)


def extract_txt(file_bytes: bytes) -> str:
    """Decode as UTF-8 with fallback to latin-1."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


def extract_html(file_bytes: bytes) -> str:
    """Basic HTML text extraction — strips tags."""
    import re
    text = extract_txt(file_bytes)
    # Remove HTML tags
    text = re.sub(r"<[^>]+>", " ", text)
    # Clean up whitespace
    text = re.sub(r"\s+", " ", text)
    return text.strip()
